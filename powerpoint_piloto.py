from pptx import Presentation
import streamlit as st
from io import BytesIO
from pptx.util import Inches
from pptx.util import Pt
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

TITLE_TEXT_SLIDE_LAYOUT = 0
BLANK_SLIDE_LAYOUT = 6
BACKGROUND_PADRAO = "images/placas.jpg"

def add_background_image(presentation, slide, image_path=BACKGROUND_PADRAO):
    # Add a picture shape
    left = top = Inches(0)
    #width = slide.width
    #height = slide.height
    pic = slide.shapes.add_picture(image_path, left, top, width=presentation.slide_width, height=presentation.slide_height)
    #https://github.com/scanny/python-pptx/issues/496
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

def set_position_object (slide, object, left, top, width, height):
    object.left = left
    object.top = top
    object.width = width
    object.height = height

def set_font(slide, paragrafo, font, size):
#https://stackoverflow.com/questions/65807855/how-change-the-font-of-a-title-of-a-slide-in-python-pptx
    title_para = slide.shapes.title.text_frame.paragraphs[paragrafo]

    title_para.font.name = font #"Comic Sans MS"
    title_para.font.size = Pt(size) #72
    title_para.font.color.rgb = RGBColor(144, 194, 38)
    #title_para.font.underline = True

def add_slide(presentation, slide_layout):
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)

    return slide

def add_title(slide, titulo):
    title = slide.shapes.title
    title.text = titulo

    return title

def add_table(slide, rows, cols, left, top, width, height):
    #https://stackoverflow.com/questions/61982333/how-to-change-default-table-style-using-pptx-python
    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    shape = slide.shapes.add_table(rows, cols, x, y, cx, cy)
    table = shape.table

    tbl =  shape._element.graphic.graphicData.tbl
    #https://github.com/scanny/python-pptx/issues/27#issuecomment-263076372
    style_id = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
    tbl[0][-1].text = style_id

    return table

def create_slide(presentation, format_slide_layout):
    st.text("Layout Format: " + str(format_slide_layout))
    slide_layout = presentation.slide_layouts[format_slide_layout]
    slide = add_slide(presentation, slide_layout)

    return slide

def apresentacao_download(presentation):
    binary_output = BytesIO()
    presentation.save(binary_output) 

    return binary_output.getvalue()


def adiciona_itens_tabela(tabela, produtos):
    #https://python-pptx.readthedocs.io/en/latest/user/table.html
    tabela.cell(0, 0).text = "Produto"

    for index_item, produto in enumerate(produtos):
        tabela.cell(index_item + 1, 0).text = produto

# remove title slide
def remove_title_slide(slide):
    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.element.getparent().remove(title_shape.element)

def remove_paragraph(slide, paragrafo):
    # Delete the paragraphs from the content placeholder
    content_shape = slide.placeholders[1]  # Adjust the index based on your slide layout
    if content_shape.has_text_frame:
        content_shape.element.getparent().remove(content_shape.element)

####################################################################################
# Cria a apresentação
####################################################################################
def create_presentation(produtos):
    presentation = Presentation()

####################################################################################
# Primeiro Slide
####################################################################################
    # Add a slide
    slide = create_slide(presentation, TITLE_TEXT_SLIDE_LAYOUT)

    # Add a title to the slide
    add_title(slide, "Sol America")
    paragrafo = 0
    font_size = 54
    set_font(slide, paragrafo, "Trebuchet MS", font_size)

    posicao_top = 10
    posicao_left = 10
    set_position_object (slide,
                         slide.shapes.title,
                         posicao_top,
                         posicao_left,
                         slide.shapes.title.width,
                         slide.shapes.title.height)

    # insert a picture from a file in slide as background
    add_background_image(presentation, slide, "images/placas.jpg")
####################################################################################
# Segundo Slide
####################################################################################
    slide = create_slide(presentation, 5)
    tabela = add_table(slide,
                       rows = 10,
                       cols = 5,
                       left = (presentation.slide_width.inches - 6.5) / 2, #3.15,#3.15,
                       top  = (presentation.slide_height.inches - 4.0) / 2, #2.20,#2.20,
                       width = 6.5,
                       height = 3.0)
    
    #slide.shapes[0].top.inches
    #slide.shapes[0].top.hight
    #slide.shapes[0].width
    #slide.shapes[0].top  = round((presentation.slide_height - Inches(1.5)) / 2)
    #slide.shapes[0].left = round((presentation.slide_width - Inches(5.5)) / 2)
    
    remove_title_slide(slide)
    remove_paragraph(slide, paragrafo = 0)

    adiciona_itens_tabela(tabela, produtos)
    
    return apresentacao_download(presentation)

def show():

    # Streamlit app
    st.title("Tales - Streamlit PowerPoint Presentation")
    st.text_input("Cliente", placeholder="Digite o nome do cliente")
    produtos = st.multiselect("Produtos", ["Produto 1", "Produto 2", "Produto 3"])

    # Insert a download button
    if st.button("Gerar aprensentação"):

        # Get the presentation
        my_presentation = create_presentation(produtos)

        st.download_button(
            label="Download Presentation",
            data=my_presentation,
            key="download_presentation",
            #on_click=lambda: st.write(my_presentation),
            help="Click to download the presentation.",
            file_name="apresentacao.pptx",
        )

show()