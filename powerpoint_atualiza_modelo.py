import streamlit as st

from pptx import Presentation
from io import BytesIO

prs = Presentation('./modelo.pptx')

# To get shapes in your slides
slides = [slide for slide in prs.slides]
shapes = []
for slide in slides:
    for shape in slide.shapes:
        shapes.append(shape)

def get_current_date():
    from datetime import date
    today = date.today()
    return today.strftime("%d/%m/%Y")

def apresentacao_download(presentation):
    binary_output = BytesIO()
    presentation.save(binary_output) 

    return binary_output.getvalue()

def replace_text(replacements: dict, shapes: list):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
      for match, replacement in replacements.items():
         if shape.has_text_frame:
            if (shape.text.find(match)) != -1:
               text_frame = shape.text_frame
               #shape.text_frame.text = shape.text_frame.text.replace(match, replacement)
               for index_paragraph, paragraph in enumerate(text_frame.paragraphs):
                  for index_run, run in enumerate(paragraph.runs):
                      shape.text_frame.paragraphs[index_paragraph].text = shape.text_frame.paragraphs[index_paragraph].text.replace(match, replacement)
         if shape.has_table:
               for row in shape.table.rows:
                  for cell in row.cells:
                     if match in cell.text:
                           new_text = cell.text.replace(match, replacement)
                           cell.text = new_text

def show():

    # Streamlit app
    st.title("Tales - Streamlit PowerPoint Presentation")
    nome_cliente = st.text_input("Cliente", placeholder="Digite o nome do cliente")
    #produtos = st.multiselect("Produtos", ["Produto 1", "Produto 2", "Produto 3"])
    potencia = st.number_input("Valor Kwp", min_value=0, step=1, value=None, placeholder="Valor Kwp")

    # Insert a download button
    if st.button("Gerar apresentação"):
      data_atual = get_current_date()
      substituicoes = {'Nome do Cliente': nome_cliente,
                       '31/12/9999': data_atual,
                       '999 Kwp': str(potencia) + ' Kwp'
                       }

      replace_text(substituicoes, shapes) 
      my_presentation = apresentacao_download(prs)

      st.download_button(
         label="Fazer Download da Apresentação",
         data=my_presentation,
         key="download_presentation",
         #on_click=lambda: st.write(my_presentation),
         help="Click para baixar a apresentação.",
         file_name="apresentacao.pptx",
      )

show()